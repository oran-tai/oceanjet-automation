"""Create a branded pptx reference template for pandoc."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Travelier brand
DARK = RGBColor(0x32, 0x37, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT = RGBColor(0x06, 0x93, 0xE3)  # Vivid cyan blue accent
HEADER_FONT = "Museo Slab"
BODY_FONT = "Raleway"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Layout 0: Title Slide ---
layout = prs.slide_layouts[0]
for ph in layout.placeholders:
    if ph.placeholder_format.idx == 0:  # Title
        ph.text_frame.paragraphs[0].font.name = HEADER_FONT
        ph.text_frame.paragraphs[0].font.size = Pt(40)
        ph.text_frame.paragraphs[0].font.color.rgb = DARK
        ph.text_frame.paragraphs[0].font.bold = True
    elif ph.placeholder_format.idx == 1:  # Subtitle
        ph.text_frame.paragraphs[0].font.name = BODY_FONT
        ph.text_frame.paragraphs[0].font.size = Pt(20)
        ph.text_frame.paragraphs[0].font.color.rgb = DARK

# --- Layout 1: Title and Content (most used) ---
layout = prs.slide_layouts[1]
for ph in layout.placeholders:
    if ph.placeholder_format.idx == 0:  # Title
        ph.text_frame.paragraphs[0].font.name = HEADER_FONT
        ph.text_frame.paragraphs[0].font.size = Pt(32)
        ph.text_frame.paragraphs[0].font.color.rgb = DARK
        ph.text_frame.paragraphs[0].font.bold = True
    elif ph.placeholder_format.idx == 1:  # Body
        ph.text_frame.paragraphs[0].font.name = BODY_FONT
        ph.text_frame.paragraphs[0].font.size = Pt(18)
        ph.text_frame.paragraphs[0].font.color.rgb = DARK
        ph.text_frame.word_wrap = True

# --- Layout 2: Section Header ---
layout = prs.slide_layouts[2]
for ph in layout.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text_frame.paragraphs[0].font.name = HEADER_FONT
        ph.text_frame.paragraphs[0].font.size = Pt(36)
        ph.text_frame.paragraphs[0].font.color.rgb = DARK
        ph.text_frame.paragraphs[0].font.bold = True

out = ".claude/skills/slides/reference-template.pptx"
prs.save(out)
print(f"Template saved to {out}")

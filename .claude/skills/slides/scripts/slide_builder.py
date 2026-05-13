"""
Reusable slide generation library for the /slides skill.

Usage:
    from slide_builder import create_presentation

    slides_data = [
        {"type": "title", "title": "My Deck", "subtitle": "Context"},
        {"type": "section", "title": "Section Name"},
        {"type": "content", "title": "Slide Title",
         "callout": "Big number or key metric",
         "bullets": [
             "Simple bullet",
             {"text": "Bullet with detail", "sub_bullets": ["Detail 1", "Detail 2"]}
         ],
         "quotes": ["\"Quote text\" — Source"],
         "table": {"headers": ["Col1", "Col2"], "rows": [["a", "b"]]}
        }
    ]
    create_presentation(slides_data, "output.pptx")
"""
import json
import math
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Rough text-height estimation for proportional body fonts (Raleway-like).
# Tuned conservatively so we under-report chars-per-line and over-report height,
# preventing visual overflow.
_CHARS_PER_INCH_BY_PT = {14: 9.5, 15: 8.8, 18: 7.4, 24: 5.5, 28: 4.7}
_LINE_HEIGHT_BY_PT = {14: 0.22, 15: 0.24, 18: 0.30, 24: 0.40, 28: 0.46}
_PADDING_BETWEEN = 0.2  # inches between content blocks


def estimate_text_height(text, font_size_pt, width_inches=11.933):
    """Estimate rendered text height in inches for a given string + font size."""
    if not text:
        return 0.0
    cpi = _CHARS_PER_INCH_BY_PT.get(font_size_pt, 7.4)
    lh = _LINE_HEIGHT_BY_PT.get(font_size_pt, 0.30)
    chars_per_line = max(1.0, width_inches * cpi)
    # Account for explicit line breaks within text
    n_lines = 0
    for line in text.split("\n"):
        n_lines += max(1, math.ceil(len(line) / chars_per_line))
    return n_lines * lh


def check_overflow(slide, slide_no, slide_height_inches=7.5):
    """Inspect a slide post-build and warn about anything that runs off the bottom."""
    warnings = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        top_in = sh.top / 914400
        h_in = sh.height / 914400
        bottom = top_in + h_in
        text = sh.text_frame.text.strip().split("\n")[0][:60]
        if bottom > slide_height_inches + 0.02:
            warnings.append(
                f"  ⚠ slide {slide_no}: shape bottom={bottom:.2f}\" exceeds {slide_height_inches}\" "
                f"(text: {text!r})"
            )
    return warnings

# Brand constants
DARK = RGBColor(0x32, 0x37, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT = RGBColor(0x07, 0x2F, 0x2F)
QUOTE_GRAY = RGBColor(0x66, 0x66, 0x66)
HEADER_FONT = "Museo Slab"
BODY_FONT = "Raleway"

# Slide dimensions
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
LEFT_MARGIN = Inches(0.7)
CONTENT_WIDTH = Inches(11.933)
TITLE_TOP = Inches(0.5)
CONTENT_TOP = Inches(1.8)


def add_textbox(slide, left, top, width, height, text, font_name, font_size,
                color=None, bold=False, italic=False, alignment=PP_ALIGN.LEFT):
    if color is None:
        color = DARK
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    return txBox


def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(1, left, top, width, Inches(0.035))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_accent_line_vertical(slide, left, top, height):
    shape = slide.shapes.add_shape(1, left, top, Inches(0.04), height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, LEFT_MARGIN, Inches(2.2), CONTENT_WIDTH, Inches(1.5),
                data["title"], HEADER_FONT, 40, DARK, bold=True, alignment=PP_ALIGN.CENTER)
    if data.get("subtitle"):
        add_textbox(slide, LEFT_MARGIN, Inches(3.8), CONTENT_WIDTH, Inches(1.0),
                    data["subtitle"], BODY_FONT, 20, DARK, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(5.5), Inches(3.5), Inches(2.333))


def add_section_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, LEFT_MARGIN, Inches(2.8), CONTENT_WIDTH, Inches(1.5),
                data["title"], HEADER_FONT, 36, DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(5.5), Inches(4.2), Inches(2.333))


def add_content_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    add_textbox(slide, LEFT_MARGIN, TITLE_TOP, CONTENT_WIDTH, Inches(0.8),
                data["title"], HEADER_FONT, 28, DARK, bold=True)
    add_accent_line(slide, LEFT_MARGIN, Inches(1.2), Inches(1.5))

    current_top = CONTENT_TOP
    current_top_inches = 1.8

    # Callout (big number / key metric) — dynamically sized
    if data.get("callout"):
        callout_h_in = max(0.5, estimate_text_height(data["callout"], 24))
        add_textbox(slide, LEFT_MARGIN, current_top, CONTENT_WIDTH, Inches(callout_h_in),
                    data["callout"], BODY_FONT, 24, ACCENT, bold=True)
        current_top_inches += callout_h_in + _PADDING_BETWEEN
        current_top = Inches(current_top_inches)

    # Bullets (with sub-bullet support) — dynamically sized
    if data.get("bullets"):
        # Pre-estimate total bullet height so the box fits the content
        est_h = 0.0
        for b in data["bullets"]:
            if isinstance(b, dict):
                main_text = b["text"]
                sub_bullets = b.get("sub_bullets", [])
            else:
                main_text = b
                sub_bullets = []
            est_h += estimate_text_height(f"•  {main_text}", 18) + 0.10
            for sb in sub_bullets:
                est_h += estimate_text_height(f"      –  {sb}", 15) + 0.05
        # Cap so the box never extends past slide bottom
        max_h = max(0.5, 7.5 - current_top_inches - 0.1)
        box_h = min(est_h + 0.2, max_h)
        txBox = slide.shapes.add_textbox(LEFT_MARGIN, current_top, CONTENT_WIDTH, Inches(box_h))
        tf = txBox.text_frame
        tf.word_wrap = True
        first = True
        for bullet in data["bullets"]:
            if isinstance(bullet, dict):
                main_text = bullet["text"]
                sub_bullets = bullet.get("sub_bullets", [])
            else:
                main_text = bullet
                sub_bullets = []

            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
                p.space_before = Pt(8)

            p.text = f"\u2022  {main_text}"
            p.font.name = BODY_FONT
            p.font.size = Pt(18)
            p.font.color.rgb = DARK

            for sb in sub_bullets:
                p2 = tf.add_paragraph()
                p2.space_before = Pt(4)
                p2.text = f"      \u2013  {sb}"
                p2.font.name = BODY_FONT
                p2.font.size = Pt(15)
                p2.font.color.rgb = DARK

        # Track vertical position — sum estimated heights for each bullet (with wrapping)
        bullets_h_in = 0.0
        for bullet in data["bullets"]:
            if isinstance(bullet, dict):
                main_text = bullet["text"]
                sub_bullets = bullet.get("sub_bullets", [])
            else:
                main_text = bullet
                sub_bullets = []
            bullets_h_in += estimate_text_height(f"•  {main_text}", 18)
            bullets_h_in += 0.10  # space_before approximation
            for sb in sub_bullets:
                bullets_h_in += estimate_text_height(f"      –  {sb}", 15)
                bullets_h_in += 0.05
        current_top_inches += bullets_h_in + _PADDING_BETWEEN
        current_top = Inches(current_top_inches)

    # Quotes
    if data.get("quotes"):
        quote_top = current_top
        quotes_h_in = 0.0
        for q in data["quotes"]:
            quotes_h_in += estimate_text_height(q, 15)
            quotes_h_in += 0.10
        add_accent_line_vertical(slide, Inches(0.9), quote_top, Inches(max(0.35, quotes_h_in)))

        txBox = slide.shapes.add_textbox(Inches(1.1), quote_top, Inches(10.5),
                                         Inches(max(0.5, quotes_h_in)))
        tf = txBox.text_frame
        tf.word_wrap = True
        first = True
        for quote in data["quotes"]:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
                p.space_before = Pt(6)
            p.text = quote
            p.font.name = BODY_FONT
            p.font.size = Pt(15)
            p.font.color.rgb = QUOTE_GRAY
            p.font.italic = True

        current_top_inches += quotes_h_in + _PADDING_BETWEEN
        current_top = Inches(current_top_inches)

    # Table
    if data.get("table"):
        tbl_data = data["table"]
        headers = tbl_data["headers"]
        rows = tbl_data["rows"]
        n_rows = len(rows) + 1
        n_cols = len(headers)

        # Tables sit below whatever came before (bullets/quotes/callout)
        table_top = current_top if (data.get("bullets") or data.get("callout") or data.get("quotes")) else CONTENT_TOP

        # Clamp to prevent going off slide
        max_top = Inches(7.5 - 0.4 * n_rows - 0.3)
        if table_top > max_top:
            table_top = max_top

        row_height = Inches(0.4)
        table_height = row_height * n_rows

        table_shape = slide.shapes.add_table(
            n_rows, n_cols, LEFT_MARGIN, table_top, CONTENT_WIDTH, table_height
        )
        table = table_shape.table

        # Style header row
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK
            for p in cell.text_frame.paragraphs:
                p.font.name = BODY_FONT
                p.font.size = Pt(14)
                p.font.color.rgb = WHITE
                p.font.bold = True

        # Style data rows
        for i, row in enumerate(rows):
            for j, value in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = value
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY if i % 2 == 1 else WHITE
                for p in cell.text_frame.paragraphs:
                    p.font.name = BODY_FONT
                    p.font.size = Pt(14)
                    p.font.color.rgb = DARK


def create_presentation(slides_data, output_path):
    """Create a branded pptx presentation from structured slide data."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for slide_data in slides_data:
        slide_type = slide_data["type"]
        if slide_type == "title":
            add_title_slide(prs, slide_data)
        elif slide_type == "section":
            add_section_slide(prs, slide_data)
        elif slide_type == "content":
            add_content_slide(prs, slide_data)

    prs.save(output_path)

    # Post-build overflow check
    all_warnings = []
    for i, slide in enumerate(prs.slides, 1):
        all_warnings.extend(check_overflow(slide, i))
    if all_warnings:
        print(f"Created {len(prs.slides)} slides \u2192 {output_path}  (with {len(all_warnings)} overflow warnings)")
        for w in all_warnings:
            print(w)
    else:
        print(f"Created {len(prs.slides)} slides \u2192 {output_path}  (no overflow detected)")


# CLI: accept JSON file as argument
if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python slide_builder.py <slides.json> <output.pptx>")
        sys.exit(1)
    with open(sys.argv[1]) as f:
        data = json.load(f)
    create_presentation(data, sys.argv[2])
